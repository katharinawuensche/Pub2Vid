#%%
from pptx import Presentation
import tempfile
import base64
import lxml.etree as etree
import pandas as pd
#%%
df = pd.read_csv("videoData.csv", delimiter="; ")

#%%
informativity, creativity, enthusiasm = 2, 2, 2
length_in_minutes = 7


inf_match = df[df["informativity"] >= informativity]
enth_match = inf_match[inf_match["enthusiasm"] >= enthusiasm]
crea_match = enth_match[enth_match["creativity"] >= creativity]
duration_match = crea_match[crea_match["duration_in_seconds"] >= (length_in_minutes-1.5)*60]
duration_match = duration_match[duration_match["duration_in_seconds"] <= (length_in_minutes+1.5)*60]
good_indices = (duration_match["informativity"] + duration_match["enthusiasm"] + duration_match["creativity"]).sort_values(ascending=False).index
duration_match.loc[good_indices[:5]]["video_id"]
















#%%
prs = Presentation("template.pptx")

#%%
# <p:transition spd="slow" advClick="0" advTm="5000"/> -> specifies the transition
""" lyt=prs.slide_layouts[13] # choosing a slide layout
slide = prs.slides.add_slide(lyt) # adding a slide
print(etree.tostring(slide._element))
for transition in slide._element.xpath("//p:transition"):
    transition.set("advTm", str(9000)) """

for slide in prs.slides:
    slide.name = "TestName"
    print(etree.tostring(slide._element))

#%%
prs = Presentation()
lyt=prs.slide_layouts[0] # choosing a slide layout
slide=prs.slides.add_slide(lyt) # adding a slide
title=slide.shapes.title # assigning a title
subtitle=slide.placeholders[1] # placeholder for subtitle
title.text="Hey,This is a Slide! How exciting!" # title
subtitle.text="Really?" # subtitle
notes = slide.notes_slide
notes.notes_text_frame.text = "foobar"
temp = tempfile.TemporaryFile()
prs.save(temp)
# %%

temp.seek(0)
print(str(base64.b64encode(bytes((temp.read()))))[:50])

# %%
temp.close()
# %%
