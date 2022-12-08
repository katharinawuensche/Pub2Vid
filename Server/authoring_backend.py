from pptx import Presentation
import pandas as pd
import tempfile
import base64
import random
import json
import datetime

sections = {
    "Title": {
        "min": 3,
        "max": 7,
        "templateIdx": 0
    }, "Introduction": {
        "min": 10,
        "max": 60,
        "templateIdx": 1
    }, "VideoOutline": {
        "min": 3,
        "max": 10,
        "templateIdx": 2
    }, "RelatedWork": {
        "min": 30,
        "max": 60,
        "templateIdx": 3
    }, "Method": {
        "min": 20,
        "max": 90,
        "templateIdx": 4
    }, "SystemArchitecture":{
        "min": 15,
        "max": 90,
        "templateIdx": 5
    }, "Demo": {
        "min": 20,
        "max": 180,
        "templateIdx": 6
    }, "Evaluation": {
        "min": 30,
        "max": 120,
        "templateIdx": 7
    }, "Results": {
        "min": 20,
        "max": 120,
        "templateIdx": 8
    }, "Reflection": {
        "min": 17,
        "max": 60,
        "templateIdx": 9
    }, "ForwardLooking": {
        "min": 10,
        "max": 40,
        "templateIdx": 10
    }, "End": {
        "min": 3,
        "max": 7,
        "templateIdx": 11
    }, "Other": {
        "min": 1,
        "max": 100,
        "templateIdx": 0
    }
}

df = pd.read_csv("videoData.csv", delimiter="; ")
videoInfo = json.load(open("preparedData.json"))

def jaccard_similarity(a, b):
    intersection = len(list(set(a).intersection(set(b))))
    union = len(set(a).union(set(b)))
    return float(intersection)/union

def get_matching_videos(d, c, e, i, v=df, n=6, top=6):
    input_ratings = [c, e, i]
    videos = v[v["duration_in_seconds"] >= d-60]
    videos = videos[videos["duration_in_seconds"] <= d+60]
    videos["weighted_sum"] = c*videos["creativity"] + e*videos["enthusiasm"] + i*videos["informativity"]
    videos = videos[videos["weighted_sum"] >= 3]
    videos["rating_sum"] = videos["creativity"] + videos["enthusiasm"] + videos["informativity"]
    videos["distance"] = 0
    for (idx, category) in enumerate(["creativity", "enthusiasm", "informativity"]):
        rel_val = videos[category]/videos["rating_sum"]
        videos["rel_"+category] = rel_val
        videos["distance"] += (rel_val - input_ratings[idx])**2
    return videos.sort_values("distance").head(top).sort_values("weighted_sum", ascending=False).head(n)

def get_video_outline(id):
    for vi in videoInfo:
        if vi["id"] == id:
            return vi["info"]["normalized_sections"]
    return []

def get_audio_elements(id):
    for vi in videoInfo:
        if vi["id"] == id:
            return vi["info"]["audioElements"]
    return []

def get_visual_elements(id):
    for vi in videoInfo:
        if vi["id"] == id:
            return vi["info"]["visualElements"]
    return []

def get_best_outline(video_ids):
    outlines = [[section["val"] for section in get_video_outline(id)] for id in video_ids]
    outline_similarities = {}
    for (idx, id) in enumerate(video_ids):
        similarity_sum = 0
        for (oidx, o2) in enumerate(outlines):
            if (oidx != idx):
                similarity_sum += jaccard_similarity(outlines[idx], o2)
        outline_similarities[id] = similarity_sum
    sorted_ids = sorted(outline_similarities.items(), key=lambda x: x[1], reverse=True)
    return get_video_outline(sorted_ids[0][0])

def get_audio_recommendations(video_ids, frac=0.5):
    all_audio_elements= [e for vid in video_ids for e in get_audio_elements(vid)]
    audio_element_count = {}
    for e in list(set(all_audio_elements)):
        audio_element_count[e] = all_audio_elements.count(e)
    return [e[0] for e in audio_element_count.items() if e[1] >= len(video_ids)*frac]

def get_visual_recommendations(video_ids, frac=0.6):
    all_visual_elements= [e for vid in video_ids for e in get_visual_elements(vid)]
    visual_element_count = {}
    for e in list(set(all_visual_elements)):
        visual_element_count[e] = all_visual_elements.count(e)
    return [e[0] for e in visual_element_count.items() if e[1] >= len(video_ids)*frac]


def generate_video_outline(duration_in_seconds=600, creativity=0, enthusiasm=0, informativity=0):
    predefinedOutlines = {
        "informativity": [
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 28 },
            { "section": "Reflection", "duration": 25 },
            { "section": "Method", "duration": 52 },
            { "section": "Evaluation", "duration": 45 },
            { "section": "Results", "duration": 53 },
            { "section": "End", "duration": 8 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 28 },
            { "section": "SystemArchitecture", "duration": 29 },
            { "section": "Method", "duration": 52 },
            { "section": "Results", "duration": 53 },
            { "section": "Evaluation", "duration": 45 },
            { "section": "End", "duration": 8 }],
            
            [{ "section": "Introduction", "duration": 28 },
            { "section": "Method", "duration": 52 },
            { "section": "Demo", "duration": 48 },
            { "section": "Results", "duration": 53 },
            { "section": "End", "duration": 8 }]
        ], 
        "enthusiasm": [
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 33 },
            { "section": "Reflection", "duration": 48 },
            { "section": "Method", "duration": 41 },
            { "section": "Evaluation", "duration": 52 },
            { "section": "Results", "duration": 53 },
            { "section": "End", "duration": 9 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 33 },
            { "section": "SystemArchitecture", "duration": 22 },
            { "section": "Method", "duration": 41 },
            { "section": "Results", "duration": 53 },
            { "section": "End", "duration": 9 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 33 },
            { "section": "SystemArchitecture", "duration": 22 },
            { "section": "Demo", "duration": 46 },
            { "section": "End", "duration": 9 }],
            
        ],
        "creativity": [
            [ { "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 28 },
            { "section": "Method", "duration": 40 },
            { "section": "Results", "duration": 50 },
            { "section": "Evaluation", "duration": 52 },
            { "section": "End", "duration": 8 }],
            [ { "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 28 },
            { "section": "Method", "duration": 40 },
            { "section": "SystemArchitecture", "duration": 22 },
            { "section": "Results", "duration": 50 },
            { "section": "End", "duration": 8 }],
           [ { "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 28 },
            { "section": "Method", "duration": 40 },
            { "section": "Demo", "duration": 44 },
            { "section": "End", "duration": 8 }],
            
        ],
        "generic": [
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 25 },
            { "section": "Method", "duration": 54 },
            { "section": "Results", "duration": 52 },
            { "section": "End", "duration": 8 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 25 },
            { "section": "Demo", "duration": 48 },
            { "section": "Method", "duration": 54 },
            { "section": "Evaluation", "duration": 52 },
            { "section": "End", "duration": 8 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "SystemArchitecture", "duration": 24 },
            { "section": "Demo", "duration": 48 },
            { "section": "End", "duration": 8 }]
        ]
    }

    #get selected outline by comparing ratio of category ratings
    selectedOutlines = []
    if creativity > (enthusiasm+informativity):
        selectedOutlines = predefinedOutlines["creativity"]
    elif enthusiasm > (creativity+informativity):
        selectedOutlines = predefinedOutlines["enthusiasm"]
    elif informativity > (enthusiasm+creativity):
        selectedOutlines = predefinedOutlines["informativity"]
    else:
        selectedOutlines = predefinedOutlines["generic"]

    for selectedOutline in selectedOutlines:
        """ for entry in selectedOutline:
            entry["duration"] = entry["duration"] * random.gauss(1.5, 0.3) """
        
        length_in_seconds = duration_in_seconds
        outlineLength = sum([entry["duration"] for entry in selectedOutline])
        scalingFactor = length_in_seconds / outlineLength

        running_seconds = 0
        for entry in selectedOutline:
            entry["val"] = entry["section"]
            entry["duration"] = entry["duration"] * scalingFactor
            entry["timeRange"] = [datetime.datetime(1970, 1, 1) + datetime.timedelta(0, running_seconds), datetime.datetime(1970, 1, 1) + datetime.timedelta(0, running_seconds + entry["duration"])]
            running_seconds += entry["duration"]
            #entry["duration"] = max(min(round(entry["duration"] * scalingFactor), sections[entry["section"]]["max"]), sections[entry["section"]]["min"])

    return selectedOutlines

def get_presentation_from_layout(outline=[]):
    prs = Presentation("template_empty.pptx")
    templateOffset = 11
    for section in outline:
        lyt=prs.slide_layouts[sections[section["section"]]["templateIdx"]+templateOffset] # choosing a slide layout
        for transition in lyt._element.xpath("//p:transition"):
            transition.set("advTm", str(section["duration"]*1000))
        slide = prs.slides.add_slide(lyt) # adding a slide
        notes = slide.notes_slide
        notes.notes_text_frame.text = str(section["notes"])
    temp = tempfile.TemporaryFile()
    prs.save(temp)
    temp.seek(0)
    res = base64.b64encode(bytes((temp.read())))
    temp.close()
    return res
    
""" 
def get_matching_videos(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    inf_match = df[df["informativity"] >= informativity]
    enth_match = inf_match[inf_match["enthusiasm"] >= enthusiasm]
    crea_match = enth_match[enth_match["creativity"] >= creativity]
    duration_match = crea_match[crea_match["duration_in_seconds"] >= (length_in_minutes-1.5)*60]
    duration_match = duration_match[duration_match["duration_in_seconds"] <= (length_in_minutes+1.5)*60]
    return duration_match """
""" 
def get_num_videos(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    crea_match = get_matching_videos(length_in_minutes, creativity, enthusiasm, informativity)
    return len(crea_match), len(df)-len(crea_match)

def get_audio_recommendations(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    crea_match = get_matching_videos(length_in_minutes, creativity, enthusiasm, informativity)
    audio_elems = [a for arr in list(crea_match["audio_elements"]) for a in str(arr).split(", ") ]
    max_audio_elems = [a for a in audio_elems if audio_elems.count(a) >= len(crea_match)*0.5]
    return list(set(max_audio_elems))

def get_visual_recommendations(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    crea_match = get_matching_videos(length_in_minutes, creativity, enthusiasm, informativity)
    visual_elems = [a for arr in list(crea_match["visual_elements"]) for a in str(arr).split(", ") ]
    max_visual_elems = [a for a in visual_elems if visual_elems.count(a) >= len(crea_match)*0.5]
    return list(set(max_visual_elems))

def get_section_recommendations(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    crea_match = get_matching_videos(length_in_minutes, creativity, enthusiasm, informativity)
    section_elems = [a for arr in list(crea_match["sections"]) for a in str(arr).split(", ") ]
    max_section_elems = [a for a in section_elems if section_elems.count(a) >= len(crea_match)*0.5]
    return list(set(max_section_elems))

def get_video_recommendations(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    crea_match = get_matching_videos(length_in_minutes, creativity, enthusiasm, informativity)
    good_indices = (crea_match["informativity"] + crea_match["enthusiasm"] + crea_match["creativity"]).sort_values(ascending=False).index
    #print(len(good_indices))
    urls = ["https://www.youtube.com/watch?v=" + id for id in list(crea_match.loc[good_indices[:5]]["video_id"])]
    print(urls)
    return urls

def get_presentation_from_length_and_sections(length_in_minutes=5, sectionNames=[]):
    #create outline
    selectedOutline = []
    for key in sections:
        if key in sectionNames:
            selectedOutline.append({"section": key, "duration": (sections[key]["max"]+sections[key]["min"])/2})

    length_in_seconds = length_in_minutes * 60
    outlineLength = sum([entry["duration"] for entry in selectedOutline])
    scalingFactor = length_in_seconds / outlineLength

    for entry in selectedOutline:
        entry["duration"] = max(min(round(entry["duration"] * scalingFactor), sections[entry["section"]]["max"]), sections[entry["section"]]["min"])
    
    return get_presentation_from_layout(selectedOutline)


def get_presentation_from_layout(outline=[]):
    prs = Presentation("template_empty.pptx")
    templateOffset = 11
    for section in outline:
        lyt=prs.slide_layouts[sections[section["section"]]["templateIdx"]+templateOffset] # choosing a slide layout
        for transition in lyt._element.xpath("//p:transition"):
            transition.set("advTm", str(section["duration"]*1000))
        slide = prs.slides.add_slide(lyt) # adding a slide
        notes = slide.notes_slide
        notes.notes_text_frame.text = str(section["notes"])
    temp = tempfile.TemporaryFile()
    prs.save(temp)
    temp.seek(0)
    res = base64.b64encode(bytes((temp.read())))
    temp.close()
    return res


def get_example_slide():
    prs = Presentation()
    lyt=prs.slide_layouts[0] # choosing a slide layout
    slide=prs.slides.add_slide(lyt) # adding a slide
    title=slide.shapes.title # assigning a title
    subtitle=slide.placeholders[1] # placeholder for subtitle
    title.text="This is a Slide! How exciting!" # title
    subtitle.text="Yay!" # subtitle
    temp = tempfile.TemporaryFile()
    prs.save(temp)
    temp.seek(0)
    res = base64.b64encode(bytes((temp.read())))
    temp.close()
    return res

def get_audio_recommendations_old(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    elements = []
    if creativity > 0.3:
        elements.append("Music")
    if enthusiasm > 0.3:
        elements.append("Music")
        elements.append("VoiceOver")
    if creativity > 0.6:
        elements.append("SoundEffects")
    if enthusiasm > 0.6:
        elements.append("SoundEffects")
    if informativity > 0.3:
        elements.append("VoiceOver")

    elements = list(set(elements))
    return elements

def get_visual_recommendations_old(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    elements = []
    if creativity > 0.3:
        elements.append("Chart")
    if enthusiasm > 0.3:
        elements.append("Chart")
        elements.append("Text")
    if creativity > 0.6:
        elements.append("Picture")
    if enthusiasm > 0.6:
        elements.append("Picture")
    if informativity > 0.3:
        elements.append("Text")

    elements = list(set(elements))
    return elements

def get_video_recommendations_old(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    return ["https://www.youtube.com/watch?v=DkF8P3qePio",
                    "https://www.youtube.com/watch?v=x4O8pojMF0w",
                    "https://www.youtube.com/watch?v=ypVWlWcR7Qk",
                    "https://www.youtube.com/watch?v=7c6oQP1u2eQ",
                    "https://www.youtube.com/watch?v=iym8fWxT9QA"]

def generate_video_outline(length_in_minutes=5, creativity=0, enthusiasm=0, informativity=0):
    predefinedOutlines = {
        "informativity": [
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 39 },
            { "section": "Method", "duration": 50 },
            { "section": "Results", "duration": 52 },
            { "section": "End", "duration": 8 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 39 },
            { "section": "SystemArchitecture", "duration": 30 },
            { "section": "Demo", "duration": 49 },
            { "section": "End", "duration": 8 }],
            [{ "section": "Introduction", "duration": 39 },
            { "section": "Method", "duration": 50 },
            { "section": "Demo", "duration": 49 },
            { "section": "End", "duration": 8 }]
        ], 
        "enthusiasm": [
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 33 },
            { "section": "Method", "duration": 30 },
            { "section": "Demo", "duration": 47 },
            { "section": "End", "duration": 9 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 33 },
            { "section": "Method", "duration": 30 },
            { "section": "SystemArchitecture", "duration": 22 },
            { "section": "Demo", "duration": 47 },
            { "section": "End", "duration": 9 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "SystemArchitecture", "duration": 22 },
            { "section": "Demo", "duration": 47 },
            { "section": "End", "duration": 9 }]
        ],
        "creativity": [
           [ { "section": "Title", "duration": 6 },
            { "section": "Introduction", "duration": 35 },
            { "section": "SystemArchitecture", "duration": 20 },
            { "section": "Demo", "duration": 48 },
            { "section": "End", "duration": 10 }],
            [ { "section": "Title", "duration": 6 },
            { "section": "Introduction", "duration": 35 },
            { "section": "Method", "duration": 28 },
            { "section": "SystemArchitecture", "duration": 20 },
            { "section": "Demo", "duration": 48 },
            { "section": "End", "duration": 10 }],
            [{ "section": "Introduction", "duration": 35 },
            { "section": "Title", "duration": 6 },
            { "section": "Demo", "duration": 48 },
            { "section": "End", "duration": 10 }]
        ],
        "generic": [
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 37 },
            { "section": "Method", "duration": 35 },
            { "section": "Demo", "duration": 48 },
            { "section": "Results", "duration": 52 },
            { "section": "End", "duration": 9 }],
            [{ "section": "Title", "duration": 5 },
            { "section": "Introduction", "duration": 37 },
            { "section": "Demo", "duration": 48 },
            { "section": "Method", "duration": 35 },
            { "section": "End", "duration": 9 }],
            [{ "section": "Introduction", "duration": 37 },
            { "section": "Method", "duration": 35 },
            { "section": "Demo", "duration": 48 },
            { "section": "Results", "duration": 52 },
            { "section": "End", "duration": 9 }]
        ]
    }

    #get selected outline by comparing ratio of category ratings
    selectedOutlines = []
    if creativity > (enthusiasm+informativity):
        selectedOutlines = predefinedOutlines["creativity"]
    elif enthusiasm > (creativity+informativity):
        selectedOutlines = predefinedOutlines["enthusiasm"]
    elif informativity > (enthusiasm+creativity):
        selectedOutlines = predefinedOutlines["informativity"]
    else:
        selectedOutlines = predefinedOutlines["generic"]

    for selectedOutline in selectedOutlines:
        for entry in selectedOutline:
            entry["duration"] = entry["duration"] * random.gauss(1.5, 0.3)
        
        length_in_seconds = length_in_minutes * 60
        outlineLength = sum([entry["duration"] for entry in selectedOutline])
        scalingFactor = length_in_seconds / outlineLength

        for entry in selectedOutline:
            entry["duration"] = max(min(round(entry["duration"] * scalingFactor), sections[entry["section"]]["max"]), sections[entry["section"]]["min"])

    return selectedOutlines
     """
